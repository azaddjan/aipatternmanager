"""
PPTX Export Service.
Generates a PowerPoint presentation matching the reference design:
dark navy theme, 10×5.625" slides, page numbers, detailed category deep-dives.
"""
import io
import os
import re
from datetime import datetime
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from services.neo4j_service import Neo4jService, BUILTIN_CATEGORIES


# ── Color Palette (matching reference) ──

BG_TITLE = RGBColor(0x0F, 0x15, 0x35)      # Title/closing slides
BG_CONTENT = RGBColor(0x1E, 0x27, 0x61)     # Content slides
BG_DIVIDER = RGBColor(0x15, 0x1D, 0x47)     # Section dividers
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xCA, 0xDC, 0xFC)     # Subtitle accent
MUTED = RGBColor(0x8E, 0x99, 0xB8)          # Body text / muted
BLUE = RGBColor(0x4A, 0x90, 0xD9)           # Primary accent
GREEN = RGBColor(0x2E, 0xCC, 0x71)          # SBB accent
PURPLE = RGBColor(0x9B, 0x59, 0xB6)         # PBC accent
ORANGE = RGBColor(0xE6, 0x7E, 0x22)         # AB/Blueprint accent
LIGHT_BLUE_2 = RGBColor(0x6B, 0xAA, 0xF0)   # Audience highlight
CARD_BG = RGBColor(0x23, 0x2D, 0x54)         # Card backgrounds
ROW_BG = RGBColor(0x1A, 0x24, 0x48)          # Darker row backgrounds
ROW_ACCENT = RGBColor(0x2A, 0x36, 0x66)      # Muted accent (L1/L2)
RED = RGBColor(0xE7, 0x4C, 0x3C)             # Warning/prohibited

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

CATEGORY_ORDER = ["core", "intg", "agt", "kr", "xcut", "pip"]

CATEGORY_INFO = {
    "core": {
        "label": "Core AI / LLM",
        "desc": "The central prompt engineering and LLM interaction capabilities — how the platform constructs, versions, and executes prompts across different runtimes and invocation paths.",
        "icon": "icon_core.png",
    },
    "intg": {
        "label": "Integration",
        "desc": "The connectivity layer bridging the platform to AI model providers, tool ecosystems, and ML serving infrastructure.",
        "icon": "icon_intg.png",
    },
    "agt": {
        "label": "Agents",
        "desc": "Autonomous AI capabilities combining reasoning, tool use, and memory for multi-step tasks.",
        "icon": "icon_agt.png",
    },
    "kr": {
        "label": "Knowledge & Retrieval",
        "desc": "How the platform ingests, indexes, and retrieves enterprise knowledge to ground LLM responses in factual context.",
        "icon": "icon_kr.png",
    },
    "xcut": {
        "label": "Cross-Cutting",
        "desc": "Governance and safety concerns spanning every category. AI guardrails delivered through integrated and standalone modes.",
        "icon": "icon_xcut.png",
    },
    "pip": {
        "label": "Platform Integration",
        "desc": "The vendor portability mechanism from PAT-002. Socket contracts, adapters, and service API contracts for integration.",
        "icon": "icon_pip.png",
    },
    "blueprint": {
        "label": "Architecture Topology",
        "desc": "Level 4 (Enterprise) patterns — foundational structural patterns that define the platform's shape.",
        "icon": "icon_topology.png",
    },
}


class PptxExportService:
    """Builds a PowerPoint deck matching the reference design."""

    def __init__(self, db: Neo4jService):
        self.db = db
        self._assets_dir = os.path.join(os.path.dirname(__file__), '..', 'static', 'pptx_assets')
        self._slide_num = 0
        self._total_slides = 0
        self._name_map: dict[str, str] = {}

    def generate_pptx(self, team_ids=None, team_names=None) -> bytes:
        """Generate a PPTX file and return raw bytes."""
        self._team_names = team_names or []
        data = self._fetch_all_data(team_ids=team_ids)

        # Build name map
        for p in data["patterns"]:
            self._name_map[p["id"]] = p.get("name", "")
        for t in data["technologies"]:
            self._name_map[t["id"]] = t.get("name", "")

        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        # Pre-calculate total slides
        grouped = defaultdict(list)
        for p in data["patterns"]:
            grouped[p.get("category", "other")].append(p)

        # Count: 10 static + category_deep_dive_divider(1) + per-category(1 each) +
        # intg deep dive(1) + topology divider(1) + topology detail(1) + boundary(1) +
        # relationships divider(1) + chain(1) + guardrail(1) + swap(1) +
        # inventory divider(1) + ABB inventory(1) + SBB inventory(1) +
        # summary divider(1) + closing(1)
        cat_count = sum(1 for c in CATEGORY_ORDER if c in grouped)
        has_blueprint = "blueprint" in grouped
        self._total_slides = (10 + 1 + cat_count +
                              (1 if "intg" in grouped else 0) +  # intg deep dive
                              (2 if has_blueprint else 0) +  # topology detail + boundary
                              1 +  # topology divider
                              1 + 3 +  # relationships divider + chain + guardrail + swap
                              1 + 2 +  # inventory divider + ABB + SBB
                              1 + 1)  # summary divider + closing
        self._slide_num = 0

        # ── Static intro slides (1-10) ──
        self._slide_01_title(prs, data)
        self._slide_02_agenda(prs, data)
        self._slide_section_divider(prs, "Why Composable", "Architecture",
                                    "Dynamic assembly over static design")
        self._slide_04_composable(prs)
        self._slide_section_divider(prs, "The ABB / SBB / PBC", "Framework",
                                    "TOGAF meets Gartner — from logical to physical to business")
        self._slide_06_three_layer(prs)
        self._slide_07_flow(prs, data)
        self._slide_08_framework_image(prs)
        self._slide_section_divider(prs, "Pattern Level", "Taxonomy",
                                    "Four levels of abstraction — two in scope for this catalogue")
        self._slide_10_levels(prs)

        # ── Category Deep Dives ──
        cat_names = [CATEGORY_INFO.get(c, {}).get("label", c) for c in CATEGORY_ORDER if c in grouped]
        self._slide_section_divider(prs, "Category", "Deep Dives",
                                    ", ".join(cat_names))

        cat_map = {c["code"]: c["label"] for c in data["categories"]}
        for cat_code in CATEGORY_ORDER:
            if cat_code not in grouped:
                continue
            pats = grouped[cat_code]
            cat_label = cat_map.get(cat_code, BUILTIN_CATEGORIES.get(cat_code, cat_code))
            self._make_category_slide(prs, cat_code, cat_label, pats)

            # Integration deep dive for Model Gateway
            if cat_code == "intg":
                abbs = [p for p in pats if p.get("type") == "ABB"]
                if abbs:
                    self._make_intg_deep_dive(prs, abbs[0], pats)

        # ── Architecture Topology ──
        self._slide_section_divider(prs, "Architecture", "Topology",
                                    "Level 4 (Enterprise) — the structure that all building blocks operate within")
        if has_blueprint:
            bp_pats = grouped["blueprint"]
            self._make_ab_detail_slide(prs, bp_pats)
            self._make_boundary_invariant_slide(prs)

        # ── Pattern Relationships ──
        self._slide_section_divider(prs, "Pattern", "Relationships",
                                    "How building blocks interrelate and depend on each other")
        self._make_dependency_chain_slide(prs)
        self._make_guardrail_mode_slide(prs, data)
        self._make_swappability_slide(prs, data)

        # ── Pattern Inventory ──
        abbs = [p for p in data["patterns"] if p.get("type") == "ABB"]
        sbbs = [p for p in data["patterns"] if p.get("type") == "SBB"]
        self._slide_section_divider(prs, "Pattern", "Inventory",
                                    f"Current catalogue: {len(grouped.get('blueprint', []))} Architecture Blueprint, {len(abbs)} ABBs, {len(sbbs)} SBBs")
        self._make_abb_inventory(prs, abbs, cat_map)
        self._make_sbb_inventory(prs, sbbs)

        # ── Summary & Closing ──
        self._slide_section_divider(prs, "Summary &", "Next Steps",
                                    "Current state, roadmap, and governance")
        self._make_closing_slide(prs, data)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ------------------------------------------------------------------
    # Data fetching
    # ------------------------------------------------------------------

    def _fetch_all_data(self, team_ids=None) -> dict:
        patterns, _ = self.db.list_patterns(limit=500, team_ids=team_ids)
        full_patterns = []
        for p in patterns:
            full = self.db.get_pattern_with_relationships(p["id"])
            if full:
                full_patterns.append(full)

        technologies, _ = self.db.list_technologies(limit=500)
        full_techs = []
        for t in technologies:
            full_t = self.db.get_technology_with_patterns(t["id"])
            if full_t:
                full_techs.append(full_t)

        pbcs = self.db.list_pbcs()
        categories = self.db.list_categories()

        return {
            "patterns": full_patterns,
            "technologies": full_techs,
            "pbcs": pbcs,
            "categories": categories,
        }

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _new_slide(self, prs, bg_color=None):
        """Create a blank slide with background and page number."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color or BG_CONTENT
        self._slide_num += 1
        # Page number
        tb = slide.shapes.add_textbox(Inches(8.5), Inches(5.2), Inches(1.2), Inches(0.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"{self._slide_num} / {self._total_slides}"
        run.font.size = Pt(9)
        run.font.color.rgb = MUTED
        run.font.name = 'Calibri'
        return slide

    def _tb(self, slide, left, top, width, height, text="", size=12, color=None,
            bold=False, align=None, font='Calibri'):
        """Add a textbox with a single styled paragraph."""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align:
            p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color or WHITE
        run.font.bold = bold
        run.font.name = font
        return tf

    def _tb_multi(self, slide, left, top, width, height, lines, size=10,
                  color=None, bold=False):
        """Add a textbox with multiple paragraphs."""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            # Clean markdown bold
            clean = re.sub(r'\*\*(.+?)\*\*', r'\1', line.strip())
            run.text = clean
            run.font.size = Pt(size)
            run.font.color.rgb = color or MUTED
            run.font.bold = bold
            run.font.name = 'Calibri'
        return tf

    def _try_image(self, slide, name, left, top, width=None, height=None):
        """Try to add an image from assets directory."""
        path = os.path.join(self._assets_dir, name)
        if os.path.exists(path):
            try:
                kwargs = {}
                if width:
                    kwargs['width'] = width
                if height:
                    kwargs['height'] = height
                slide.shapes.add_picture(path, left, top, **kwargs)
                return True
            except Exception:
                pass
        return False

    def _rect(self, slide, left, top, width, height, fill_color):
        """Add a filled rectangle with no border (used for card backgrounds and accent lines)."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()  # No border
        return shape

    def _get_functionality_bullets(self, pattern):
        """Extract short functionality bullets from a pattern."""
        func = pattern.get("functionality") or pattern.get("specific_functionality") or ""
        if not func:
            return []
        lines = []
        for line in func.split("\n"):
            stripped = line.strip()
            if stripped.startswith("- ") or stripped.startswith("* "):
                clean = re.sub(r'\*\*(.+?)\*\*', r'\1', stripped[2:])
                lines.append(clean)
            elif stripped and len(lines) < 8:
                clean = re.sub(r'\*\*(.+?)\*\*', r'\1', stripped)
                if len(clean) < 80:
                    lines.append(clean)
        return lines[:8]

    def _get_short_desc(self, pattern):
        """Get a short one-line description from a pattern."""
        # Prefer the explicit description field first
        desc = pattern.get("description", "")
        if desc:
            first = desc.split("\n")[0].strip()
            clean = re.sub(r'\*\*(.+?)\*\*', r'\1', first)
            return clean[:80]
        for field in ["specific_functionality", "functionality", "intent"]:
            val = pattern.get(field, "")
            if val:
                first = val.split("\n")[0].strip()
                if first.startswith("- "):
                    first = first[2:]
                clean = re.sub(r'\*\*(.+?)\*\*', r'\1', first)
                return clean[:80]
        return ""

    # ------------------------------------------------------------------
    # Slide 1: Title
    # ------------------------------------------------------------------

    def _slide_01_title(self, prs, data):
        slide = self._new_slide(prs, BG_TITLE)

        # Top banner
        self._rect(slide, Inches(0), Inches(0), Inches(10), Inches(2.2), BG_DIVIDER)
        # Blue accent line
        self._rect(slide, Inches(0), Inches(2.15), Inches(10), Inches(0.04), BLUE)

        # Logo icon
        self._try_image(slide, "icon_abb.png", Inches(0.7), Inches(0.5), Inches(0.6), Inches(0.6))

        # Title
        self._tb(slide, Inches(0.7), Inches(1.1), Inches(8.5), Inches(0.8),
                 "Architecture Patterns", size=38, bold=True)

        # Subtitle
        self._tb(slide, Inches(0.7), Inches(2.6), Inches(8.0), Inches(0.5),
                 "A Composable Enterprise Framework", size=18, color=LIGHT_BLUE)

        # Description
        self._tb_multi(slide, Inches(0.7), Inches(3.2), Inches(8.0), Inches(0.6), [
            "Combining TOGAF Architecture & Solution Building Blocks",
            "with Gartner Packaged Business Capabilities",
        ], size=13, color=MUTED)

        # Stats summary
        patterns = data["patterns"]
        ab_count = sum(1 for p in patterns if p.get("type") == "AB")
        abb_count = sum(1 for p in patterns if p.get("type") == "ABB")
        sbb_count = sum(1 for p in patterns if p.get("type") == "SBB")
        tech_count = len(data["technologies"])
        stats_text = f"{ab_count} Blueprints  |  {abb_count} ABBs  |  {sbb_count} SBBs  |  {tech_count} Technologies"
        self._tb(slide, Inches(0.7), Inches(4.0), Inches(8.0), Inches(0.3),
                 stats_text, size=11, color=MUTED)

        # Export date with timezone
        from zoneinfo import ZoneInfo
        est_now = datetime.now(ZoneInfo("America/New_York"))
        export_date = est_now.strftime("%B %d, %Y  %I:%M %p EST")
        self._tb(slide, Inches(0.7), Inches(4.35), Inches(5.0), Inches(0.3),
                 export_date, size=11, color=MUTED)

    # ------------------------------------------------------------------
    # Slide 2: Agenda
    # ------------------------------------------------------------------

    def _slide_02_agenda(self, prs, data):
        slide = self._new_slide(prs, BG_CONTENT)
        self._tb(slide, Inches(0.7), Inches(0.4), Inches(8.0), Inches(0.7),
                 "Agenda", size=30, bold=True)

        items = [
            ("Why Composable Architecture", "The case for dynamic assembly over static design"),
            ("The PBC / ABB / SBB Framework", "TOGAF meets Gartner — from business to logical to physical"),
            ("Pattern Level Taxonomy", "Four levels of abstraction, two in scope"),
            ("Category Deep Dives (L3)", "Core, Integration, Agents, Knowledge, Cross-Cutting, Platform Integration"),
            ("Architecture Topology (L4)", "The Segmented Platform — structure, not a building block"),
            ("Pattern Relationships", "How building blocks interrelate and depend on each other"),
            ("Summary & Next Steps", "Current inventory and roadmap"),
        ]

        y = Inches(1.5)
        for i, (title, desc) in enumerate(items):
            # Number
            self._tb(slide, Inches(0.7), y, Inches(0.5), Inches(0.4),
                     f"{i+1:02d}", size=14, bold=True, color=BLUE)
            # Title
            self._tb(slide, Inches(1.4), y, Inches(4.5), Inches(0.2),
                     title, size=14, bold=True)
            # Description
            self._tb(slide, Inches(1.4), y + Inches(0.25), Inches(7.0), Inches(0.2),
                     desc, size=11, color=MUTED)
            y += Inches(0.55)

    # ------------------------------------------------------------------
    # Section Divider
    # ------------------------------------------------------------------

    def _slide_section_divider(self, prs, line1, line2, subtitle=""):
        slide = self._new_slide(prs, BG_DIVIDER)
        self._tb_multi(slide, Inches(0.8), Inches(1.8), Inches(8.5), Inches(1.0),
                       [line1, line2], size=36, color=WHITE, bold=True)
        if subtitle:
            self._tb(slide, Inches(0.8), Inches(2.9), Inches(8.0), Inches(0.6),
                     subtitle, size=14, color=MUTED)

    # ------------------------------------------------------------------
    # Slide 4: Composable Imperative
    # ------------------------------------------------------------------

    def _slide_04_composable(self, prs):
        slide = self._new_slide(prs, BG_CONTENT)
        self._tb(slide, Inches(0.7), Inches(0.35), Inches(8.0), Inches(0.6),
                 "The Composable Imperative", size=26, bold=True)
        self._tb(slide, Inches(0.7), Inches(1.0), Inches(8.5), Inches(0.6),
                 "The AI landscape evolves too rapidly for static blueprints. Composable architecture shifts from static design to dynamic assembly.",
                 size=12, color=MUTED)

        pillars = [
            (0.5, "Technology\nPortability",
             "Well-defined ABB contracts enable swapping between self-hosted and SaaS implementations without architectural disruption.",
             "icon_tech_portability.png"),
            (3.55, "Governance\nWithout Rigidity",
             "Governance policies attach to ABBs, not SBBs — consistent controls without constraining implementation choices.",
             "icon_governance.png"),
            (6.6, "Incremental\nAdoption",
             "Teams implement one building block at a time while maintaining architectural coherence across the enterprise.",
             "icon_adoption.png"),
        ]

        for x, title, desc, icon in pillars:
            # Card background
            self._rect(slide, Inches(x), Inches(1.85), Inches(2.8), Inches(3.2), CARD_BG)
            # Icon
            icon_x = x + 1.1
            self._try_image(slide, icon, Inches(icon_x), Inches(2.1), Inches(0.55), Inches(0.55))
            # Title
            self._tb_multi(slide, Inches(x + 0.2), Inches(2.8), Inches(2.4), Inches(0.7),
                           title.split("\n"), size=15, color=WHITE, bold=True)
            # Desc
            self._tb(slide, Inches(x + 0.2), Inches(3.55), Inches(2.4), Inches(1.2),
                     desc, size=11, color=MUTED)

    # ------------------------------------------------------------------
    # Slide 6: Three-Layer Framework
    # ------------------------------------------------------------------

    def _slide_06_three_layer(self, prs):
        slide = self._new_slide(prs, BG_CONTENT)
        self._tb(slide, Inches(0.7), Inches(0.35), Inches(8.0), Inches(0.6),
                 "Three-Layer Framework", size=26, bold=True)

        layers = [
            (1.0, "PBC", "Business / Why", "Packaged Business Capability",
             "Business-consumable services exposed via APIs, bundling one or more ABB/SBB combinations into a reusable capability recognizable by business users.",
             PURPLE, "icon_pbc.png"),
            (2.4, "ABB", "Logical / What", "Architecture Building Block",
             "Vendor-neutral specifications defining required capabilities, interfaces, and constraints. ABBs define what an architecture requires without mandating how.",
             BLUE, "icon_abb.png"),
            (3.8, "SBB", "Physical / How", "Solution Building Block",
             "Concrete implementations fulfilling ABB contracts. SBBs are product-aware and vendor-specific. Multiple SBBs can satisfy the same ABB contract.",
             GREEN, "icon_sbb.png"),
        ]

        for y, abbrev, role, full_name, desc, color, icon in layers:
            # Card background
            self._rect(slide, Inches(0.5), Inches(y), Inches(9.0), Inches(1.2), CARD_BG)
            # Left accent line
            self._rect(slide, Inches(0.5), Inches(y), Inches(0.07), Inches(1.2), color)
            # Icon
            self._try_image(slide, icon, Inches(0.85), Inches(y + 0.3), Inches(0.5), Inches(0.5))
            # Abbreviation
            self._tb(slide, Inches(1.6), Inches(y + 0.12), Inches(1.2), Inches(0.35),
                     abbrev, size=16, bold=True, color=color)
            # Role
            self._tb(slide, Inches(2.7), Inches(y + 0.15), Inches(2.0), Inches(0.3),
                     role, size=10, color=MUTED)
            # Full name
            self._tb(slide, Inches(1.6), Inches(y + 0.45), Inches(7.0), Inches(0.3),
                     full_name, size=13, bold=True)
            # Description
            self._tb(slide, Inches(1.6), Inches(y + 0.75), Inches(7.5), Inches(0.4),
                     desc, size=11, color=MUTED)

    # ------------------------------------------------------------------
    # Slide 7: How It Flows Together
    # ------------------------------------------------------------------

    def _slide_07_flow(self, prs, data):
        slide = self._new_slide(prs, BG_CONTENT)
        self._tb(slide, Inches(0.7), Inches(0.35), Inches(8.0), Inches(0.6),
                 "How It Flows Together", size=26, bold=True)
        self._tb(slide, Inches(0.7), Inches(1.0), Inches(8.5), Inches(0.5),
                 "ABB contracts define stable interfaces while allowing SBB implementations to evolve. PBCs compose ABBs into business-facing capabilities.",
                 size=12, color=MUTED)

        # Flow: PBC → ABB → SBB with card backgrounds
        flow = [
            (0.8, "PBC", "packages for business", PURPLE),
            (3.8, "ABB", "defines the contract", BLUE),
            (6.8, "SBB", "fulfills it", GREEN),
        ]
        for x, abbrev, desc, color in flow:
            # Card background
            self._rect(slide, Inches(x), Inches(2.0), Inches(2.5), Inches(1.6), CARD_BG)
            # Top accent line
            self._rect(slide, Inches(x), Inches(2.0), Inches(2.5), Inches(0.06), color)
            # Text
            self._tb(slide, Inches(x), Inches(2.3), Inches(2.5), Inches(0.5),
                     abbrev, size=22, bold=True, color=color)
            self._tb(slide, Inches(x), Inches(2.85), Inches(2.5), Inches(0.5),
                     desc, size=12, color=MUTED)

        # Arrows
        self._tb(slide, Inches(3.3), Inches(2.3), Inches(1.3), Inches(0.8),
                 "→", size=30, color=BLUE)
        self._tb(slide, Inches(6.3), Inches(2.3), Inches(1.3), Inches(0.8),
                 "→", size=30, color=BLUE)

        # Example
        self._tb(slide, Inches(0.7), Inches(4.0), Inches(8.0), Inches(0.35),
                 "Example: Prompt Engineering", size=13, bold=True)
        self._tb(slide, Inches(0.7), Inches(4.35), Inches(9.0), Inches(0.5),
                 "PBC: Natural Language Interaction  →  ABB-CORE-001 (Prompt Engineering)  →  SBB-CORE-001 (Direct EKS) / SBB-CORE-002 (Gateway)",
                 size=11, color=MUTED)

    # ------------------------------------------------------------------
    # Slide 8: Framework Image
    # ------------------------------------------------------------------

    def _slide_08_framework_image(self, prs):
        slide = self._new_slide(prs, BG_CONTENT)
        self._tb(slide, Inches(0.7), Inches(0.3), Inches(8.0), Inches(0.6),
                 "The Unified Framework", size=26, bold=True)
        self._tb(slide, Inches(0.7), Inches(0.8), Inches(8.5), Inches(0.3),
                 "PBC packages business capabilities, ABB defines vendor-neutral contracts, SBB provides concrete implementations",
                 size=11, color=MUTED)

        img_path = os.path.join(os.path.dirname(__file__), '..', 'static', 'framework_diagram.png')
        if os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.5))
            except Exception:
                pass

        self._tb(slide, Inches(0.7), Inches(4.3), Inches(8.5), Inches(0.3),
                 "PBC (Why)  →  ABB (What)  →  SBB (How)", size=12, bold=True, color=BLUE)

    # ------------------------------------------------------------------
    # Slide 10: Pattern Levels
    # ------------------------------------------------------------------

    def _slide_10_levels(self, prs):
        slide = self._new_slide(prs, BG_CONTENT)
        self._tb(slide, Inches(0.7), Inches(0.35), Inches(8.0), Inches(0.6),
                 "Pattern Levels", size=26, bold=True)
        self._tb(slide, Inches(0.7), Inches(0.9), Inches(8.0), Inches(0.35),
                 "This catalogue focuses on Level 3 (Architectural) and Level 4 (Enterprise) patterns.",
                 size=12, color=MUTED)

        levels = [
            ("L1", "Implementation", "Code & algorithm level — prompt templates, chunking strategies, embedding techniques", "Developers", False),
            ("L2", "Design", "Component level — ReAct Loop, Router Agent, Chain of Thought, Basic RAG", "Technical Leads", False),
            ("L3", "Architectural", "System level — RAG Architecture, Agent Deployment, Human-in-the-Loop, Model Gateway", "Solution Architects", True),
            ("L4", "Enterprise Patterns & Topology", "Architecture Topology, Segmented Platform, Multi-Account Deployment, Capability Packaging", "Enterprise Architects", True),
        ]

        y_val = 1.5
        for level, title, desc, audience, in_scope in levels:
            y = Inches(y_val)
            # Row background
            row_bg = CARD_BG if in_scope else ROW_BG
            self._rect(slide, Inches(0.5), y, Inches(9.0), Inches(0.8), row_bg)
            # Left accent line
            accent = BLUE if in_scope else ROW_ACCENT
            self._rect(slide, Inches(0.5), y, Inches(0.06), Inches(0.8), accent)

            lc = BLUE if in_scope else MUTED
            tc = WHITE if in_scope else MUTED
            self._tb(slide, Inches(0.8), y + Inches(0.1), Inches(0.6), Inches(0.35),
                     level, size=16, bold=True, color=lc)
            self._tb(slide, Inches(1.5), y + Inches(0.07), Inches(3.5), Inches(0.3),
                     title, size=14, bold=True, color=tc)
            self._tb(slide, Inches(1.5), y + Inches(0.38), Inches(6.5), Inches(0.3),
                     desc, size=10, color=MUTED)
            aud_color = LIGHT_BLUE_2 if in_scope else MUTED
            self._tb(slide, Inches(7.5), y + Inches(0.07), Inches(1.8), Inches(0.3),
                     audience, size=10, color=aud_color)
            if in_scope:
                self._tb(slide, Inches(8.2), y + Inches(0.45), Inches(1.0), Inches(0.25),
                         "IN SCOPE", size=8, bold=True, color=BLUE)
            y_val += 0.95

    # ------------------------------------------------------------------
    # Category Deep Dive
    # ------------------------------------------------------------------

    def _make_category_slide(self, prs, cat_code, cat_label, patterns):
        slide = self._new_slide(prs, BG_CONTENT)
        info = CATEGORY_INFO.get(cat_code, {})

        # Icon + title
        self._try_image(slide, info.get("icon", ""), Inches(0.7), Inches(0.3), Inches(0.4), Inches(0.4))
        self._tb(slide, Inches(1.25), Inches(0.3), Inches(5.0), Inches(0.5),
                 info.get("label", cat_label), size=24, bold=True)

        # Description
        self._tb(slide, Inches(0.7), Inches(0.9), Inches(8.5), Inches(0.5),
                 info.get("desc", ""), size=11, color=MUTED)

        abbs = sorted([p for p in patterns if p.get("type") == "ABB"], key=lambda p: p["id"])
        sbbs = sorted([p for p in patterns if p.get("type") == "SBB"], key=lambda p: p["id"])

        # Calculate available space: top of cards area (1.6) to bottom safe zone (5.1)
        CARDS_TOP = 1.6
        CARDS_BOTTOM = 5.1
        available_height = CARDS_BOTTOM - CARDS_TOP

        # ABB card background (left side)
        self._rect(slide, Inches(0.5), Inches(CARDS_TOP), Inches(4.3), Inches(available_height), CARD_BG)
        # Blue top accent line for ABB section
        self._rect(slide, Inches(0.5), Inches(CARDS_TOP), Inches(4.3), Inches(0.05), BLUE)

        # Determine max bullets per ABB based on how many ABBs we have
        max_bullets = 3 if len(abbs) <= 2 else 2 if len(abbs) <= 3 else 1
        MAX_BULLET_LEN = 55  # Truncate long bullets to avoid wrapping

        # Left side: ABBs with functionality bullets
        y = Inches(CARDS_TOP + 0.15)
        for abb in abbs:
            if y > Inches(CARDS_BOTTOM - 0.4):
                break
            # ABB badge + name on same line
            abb_team = abb.get("team_name", "")
            name_line = f"{abb['id']} — {abb.get('name', '')}"
            self._tb(slide, Inches(0.7), y, Inches(0.6), Inches(0.22),
                     "ABB", size=9, bold=True, color=BLUE)
            self._tb(slide, Inches(1.3), y, Inches(3.4), Inches(0.22),
                     name_line, size=11, bold=True)
            y += Inches(0.24)
            # Team badge inline (subtle, on its own line)
            if abb_team:
                self._tb(slide, Inches(1.3), y, Inches(3.0), Inches(0.15),
                         abb_team, size=7, color=LIGHT_BLUE_2)
                y += Inches(0.15)
            # Functionality bullets (limited + truncated)
            bullets = self._get_functionality_bullets(abb)[:max_bullets]
            if bullets:
                short_bullets = [
                    f"• {b[:MAX_BULLET_LEN]}…" if len(b) > MAX_BULLET_LEN else f"• {b}"
                    for b in bullets
                ]
                self._tb_multi(slide, Inches(0.9), y, Inches(3.6), Inches(len(short_bullets) * 0.16 + 0.05),
                               short_bullets, size=8, color=MUTED)
                y += Inches(len(short_bullets) * 0.16 + 0.08)
            y += Inches(0.12)

        # Right side: SBBs — calculate card height based on count
        sbb_count = len(sbbs)
        card_gap = 0.08
        if sbb_count <= 3:
            card_height = min(1.0, (available_height - card_gap * (sbb_count - 1)) / max(sbb_count, 1))
        else:
            card_height = (available_height - card_gap * (sbb_count - 1)) / sbb_count
        card_height = max(0.55, min(1.0, card_height))  # clamp between 0.55 and 1.0

        sbb_y = CARDS_TOP
        for sbb in sbbs:
            if sbb_y + card_height > CARDS_BOTTOM + 0.05:
                break
            # SBB card background
            self._rect(slide, Inches(5.2), Inches(sbb_y), Inches(4.3), Inches(card_height), CARD_BG)
            # Green left accent line
            self._rect(slide, Inches(5.2), Inches(sbb_y), Inches(0.06), Inches(card_height), GREEN)
            # SBB ID
            self._tb(slide, Inches(5.5), Inches(sbb_y + 0.06), Inches(2.0), Inches(0.18),
                     sbb["id"], size=8, bold=True, color=GREEN)
            # Name
            self._tb(slide, Inches(5.5), Inches(sbb_y + 0.24), Inches(3.7), Inches(0.2),
                     sbb.get("name", ""), size=10, bold=True)
            # Team name (right-aligned, subtle)
            sbb_team = sbb.get("team_name")
            if sbb_team:
                self._tb(slide, Inches(7.5), Inches(sbb_y + 0.06), Inches(1.8), Inches(0.15),
                         sbb_team, size=7, color=LIGHT_BLUE_2,
                         align=PP_ALIGN.RIGHT)
            # Short description (only if card is tall enough)
            if card_height >= 0.65:
                desc = self._get_short_desc(sbb)
                if desc:
                    self._tb(slide, Inches(5.5), Inches(sbb_y + 0.42), Inches(3.7), Inches(0.2),
                             desc[:60] + ("…" if len(desc) > 60 else ""), size=8, color=MUTED)
            sbb_y += card_height + card_gap

    # ------------------------------------------------------------------
    # Integration Deep Dive (Model Gateway)
    # ------------------------------------------------------------------

    def _make_intg_deep_dive(self, prs, first_abb, pats):
        slide = self._new_slide(prs, BG_CONTENT)
        abb_name = first_abb.get("name", "Model Gateway")
        self._tb(slide, Inches(0.7), Inches(0.35), Inches(8.0), Inches(0.6),
                 f"Deep Dive: {abb_name}", size=24, bold=True)
        self._tb(slide, Inches(0.7), Inches(0.85), Inches(8.5), Inches(0.35),
                 f"{first_abb['id']} — The single enforcement point for all AI/ML model invocations",
                 size=11, color=MUTED)

        # Left card: Capabilities
        self._rect(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(3.6), CARD_BG)
        self._tb(slide, Inches(0.7), Inches(1.6), Inches(3.0), Inches(0.3),
                 "Capabilities", size=13, bold=True, color=BLUE)
        bullets = self._get_functionality_bullets(first_abb)
        if bullets:
            self._tb_multi(slide, Inches(0.7), Inches(2.0), Inches(3.8), Inches(2.8),
                           bullets, size=10, color=MUTED)

        # Right card: SBB Implementations
        self._rect(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(3.6), CARD_BG)
        self._tb(slide, Inches(5.4), Inches(1.6), Inches(3.0), Inches(0.3),
                 "SBB Implementations", size=13, bold=True, color=GREEN)

        sbbs = sorted([p for p in pats if p.get("type") == "SBB"], key=lambda p: p["id"])
        y = Inches(2.1)
        for sbb in sbbs[:5]:
            self._tb(slide, Inches(5.4), y, Inches(2.0), Inches(0.2),
                     sbb["id"], size=9, bold=True, color=GREEN)
            self._tb(slide, Inches(5.4), y + Inches(0.18), Inches(3.5), Inches(0.2),
                     sbb.get("name", ""), size=11, bold=True)
            desc = self._get_short_desc(sbb)
            if desc:
                self._tb(slide, Inches(5.4), y + Inches(0.38), Inches(3.8), Inches(0.2),
                         desc, size=9, color=MUTED)
            y += Inches(0.7)

    # ------------------------------------------------------------------
    # Architecture Topology Detail
    # ------------------------------------------------------------------

    def _make_ab_detail_slide(self, prs, bp_pats):
        slide = self._new_slide(prs, BG_CONTENT)
        info = CATEGORY_INFO.get("blueprint", {})
        self._try_image(slide, info.get("icon", ""), Inches(0.7), Inches(0.3), Inches(0.4), Inches(0.4))
        self._tb(slide, Inches(1.25), Inches(0.3), Inches(5.0), Inches(0.5),
                 "Architecture Topology", size=24, bold=True)
        self._tb(slide, Inches(0.7), Inches(0.9), Inches(8.5), Inches(0.5),
                 "Level 4 (Enterprise) patterns — foundational structural patterns that define the platform's shape.",
                 size=11, color=MUTED)

        if bp_pats:
            ab = bp_pats[0]

            # Main card background with orange accent
            self._rect(slide, Inches(0.5), Inches(1.65), Inches(9.0), Inches(3.5), CARD_BG)
            self._rect(slide, Inches(0.5), Inches(1.65), Inches(0.07), Inches(3.5), ORANGE)

            self._try_image(slide, "icon_topology.png", Inches(0.85), Inches(1.85), Inches(0.45), Inches(0.45))
            self._tb(slide, Inches(1.5), Inches(1.8), Inches(2.0), Inches(0.35),
                     ab["id"], size=11, bold=True, color=ORANGE)
            self._tb(slide, Inches(1.5), Inches(2.1), Inches(5.0), Inches(0.35),
                     ab.get("name", ""), size=16, bold=True)

            # Try to extract planes from structural elements
            struct = ab.get("structural_elements", "")
            planes = []
            if struct:
                for line in struct.split("\n"):
                    stripped = line.strip()
                    if "Plane" in stripped or "plane" in stripped:
                        clean = re.sub(r'\*\*(.+?)\*\*', r'\1', stripped)
                        if clean.startswith("- "):
                            clean = clean[2:]
                        planes.append(clean)

            if not planes:
                planes = [
                    "Experience Plane — How teams consume the platform — developer portal, SDKs, self-service",
                    "Control Plane — How the platform governs — identity, policy/guardrails, agent lifecycle",
                    "Execution Plane — How AI workloads run — model inference, retrieval, agent runtime",
                ]

            # Plane accent colors
            plane_colors = [RGBColor(0x34, 0x98, 0xDB), ORANGE, GREEN]  # Blue, Orange, Green

            y_val = 2.7
            for i, plane in enumerate(planes[:4]):
                y = Inches(y_val)
                parts = plane.split("—") if "—" in plane else plane.split(" – ") if " – " in plane else [plane]
                title = parts[0].strip()
                desc = " — ".join(p.strip() for p in parts[1:]) if len(parts) > 1 else ""

                # Plane row background
                self._rect(slide, Inches(1.0), y, Inches(8.2), Inches(0.55), ROW_BG)
                accent = plane_colors[i] if i < len(plane_colors) else BLUE
                self._rect(slide, Inches(1.0), y, Inches(0.06), Inches(0.55), accent)

                self._tb(slide, Inches(1.3), y + Inches(0.02), Inches(2.5), Inches(0.25),
                         title, size=12, bold=True)
                if desc:
                    self._tb(slide, Inches(1.3), y + Inches(0.26), Inches(7.5), Inches(0.25),
                             desc, size=10, color=MUTED)
                y_val += 0.65

    # ------------------------------------------------------------------
    # Boundary Invariant
    # ------------------------------------------------------------------

    def _make_boundary_invariant_slide(self, prs):
        slide = self._new_slide(prs, BG_CONTENT)
        self._tb(slide, Inches(0.7), Inches(0.35), Inches(8.0), Inches(0.6),
                 "The Boundary Invariant", size=26, bold=True)

        # Main invariant card with red top accent
        self._rect(slide, Inches(1.0), Inches(1.3), Inches(8.0), Inches(1.0), CARD_BG)
        self._rect(slide, Inches(1.0), Inches(1.3), Inches(8.0), Inches(0.05), RED)
        self._tb(slide, Inches(1.2), Inches(1.45), Inches(7.5), Inches(0.7),
                 "Experience never bypasses Control to reach Execution.",
                 size=18, bold=True)

        boundaries = [
            ("Experience ↔ Control",
             "All consumer requests flow through Control. Authentication, agent lifecycle, rate limiting enforced before any execution.",
             BLUE),
            ("Control ↔ Execution",
             "All LLM calls route through Model Gateway. All tool calls route through Tool Gateway. No direct provider access.",
             BLUE),
            ("Experience ↔ Execution",
             "PROHIBITED. No direct contract exists. All interactions are mediated by the Control Plane.",
             RED),
        ]

        y_val = 2.8
        for title, desc, accent in boundaries:
            y = Inches(y_val)
            # Row background
            self._rect(slide, Inches(1.0), y, Inches(8.0), Inches(0.7), ROW_BG)
            # Left accent
            self._rect(slide, Inches(1.0), y, Inches(0.06), Inches(0.7), accent)
            self._tb(slide, Inches(1.3), y + Inches(0.05), Inches(3.0), Inches(0.25),
                     title, size=12, bold=True)
            self._tb(slide, Inches(1.3), y + Inches(0.32), Inches(7.3), Inches(0.3),
                     desc, size=10, color=MUTED)
            y_val += 0.85

    # ------------------------------------------------------------------
    # Dependency Chain
    # ------------------------------------------------------------------

    def _make_dependency_chain_slide(self, prs):
        slide = self._new_slide(prs, BG_CONTENT)
        self._tb(slide, Inches(0.7), Inches(0.35), Inches(8.0), Inches(0.6),
                 "Core Dependency Chain", size=24, bold=True)
        self._tb(slide, Inches(0.7), Inches(0.85), Inches(8.5), Inches(0.35),
                 "Every AI workload follows this dependency path from prompt to model to guardrails to knowledge.",
                 size=11, color=MUTED)

        # 4-box flow with card backgrounds
        boxes = [
            (0.3, "Prompt\nEngineering", "ABB-CORE-001"),
            (2.7, "Model\nGateway", "ABB-INTG-001"),
            (5.1, "AI\nGuardrails", "ABB-XCUT-001"),
            (7.5, "Knowledge &\nRetrieval", "ABB-KR-001"),
        ]

        for x, title, ref_id in boxes:
            # Card background
            self._rect(slide, Inches(x), Inches(1.5), Inches(2.1), Inches(1.5), CARD_BG)
            self._tb_multi(slide, Inches(x), Inches(1.65), Inches(2.1), Inches(0.7),
                           title.split("\n"), size=13, color=WHITE, bold=True)
            self._tb(slide, Inches(x), Inches(2.45), Inches(2.1), Inches(0.3),
                     ref_id, size=9, color=BLUE)

        # Arrows
        for x in [2.35, 4.75, 7.15]:
            self._tb(slide, Inches(x), Inches(1.7), Inches(0.7), Inches(1.0),
                     "➔", size=26, color=BLUE)

        # Key relationships
        self._tb(slide, Inches(0.7), Inches(3.4), Inches(8.0), Inches(0.3),
                 "Key Relationships", size=14, bold=True)

        rels = [
            "ABB-CORE-001 depends on ABB-INTG-001 — all prompt execution requires LLM access via the gateway",
            "ABB-INTG-001 depends on ABB-PIP-001/002 — vendor adapters implement the socket contract",
            "ABB-AGT-001 depends on ABB-CORE-001, ABB-INTG-001, ABB-INTG-002 — agents combine prompts, models, tools",
            "ABB-KR-001 depends on ABB-INTG-001 — embedding generation routes through the model gateway",
            "ABB-XCUT-001 is consumed by all Core and Agent building blocks — guardrails are cross-cutting",
        ]
        self._tb_multi(slide, Inches(0.7), Inches(3.8), Inches(8.5), Inches(1.6),
                       rels, size=10, color=MUTED)

    # ------------------------------------------------------------------
    # Guardrail Mode Selection
    # ------------------------------------------------------------------

    def _make_guardrail_mode_slide(self, prs, data):
        slide = self._new_slide(prs, BG_CONTENT)
        self._tb(slide, Inches(0.7), Inches(0.35), Inches(8.0), Inches(0.6),
                 "Guardrail Mode Selection", size=24, bold=True)
        self._tb(slide, Inches(0.7), Inches(0.85), Inches(8.5), Inches(0.35),
                 "The invocation path determines which guardrail mode is available. This is a key architectural decision.",
                 size=11, color=MUTED)

        # Left card: Direct Bedrock Path
        self._rect(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(3.5), CARD_BG)
        self._rect(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(0.05), GREEN)
        self._tb(slide, Inches(0.7), Inches(1.65), Inches(3.0), Inches(0.35),
                 "Direct Bedrock Path", size=15, bold=True)
        self._tb(slide, Inches(0.7), Inches(2.0), Inches(3.5), Inches(0.25),
                 "SBB-CORE-001 / SBB-CORE-003", size=10, color=GREEN)
        self._tb_multi(slide, Inches(0.7), Inches(2.4), Inches(3.8), Inches(1.5), [
            "App calls Bedrock InvokeModel directly",
            "Guardrails enforced inline during call",
            "Single API call = inference + safety",
            "Lower latency, simpler architecture",
        ], size=10, color=MUTED)
        self._tb(slide, Inches(0.7), Inches(4.2), Inches(3.8), Inches(0.3),
                 "Uses SBB-XCUT-001 (Integrated)", size=10, bold=True, color=BLUE)

        # Right card: Gateway-Routed Path
        self._rect(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(3.5), CARD_BG)
        self._rect(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(0.05), PURPLE)
        self._tb(slide, Inches(5.4), Inches(1.65), Inches(3.0), Inches(0.35),
                 "Gateway-Routed Path", size=15, bold=True)
        self._tb(slide, Inches(5.4), Inches(2.0), Inches(3.5), Inches(0.25),
                 "SBB-CORE-002 via LiteLLM", size=10, color=PURPLE)
        self._tb_multi(slide, Inches(5.4), Inches(2.4), Inches(3.8), Inches(1.5), [
            "App calls LiteLLM gateway",
            "LiteLLM routes to any provider",
            "Inline Bedrock guardrails unavailable",
            "Requires separate ApplyGuardrail API call",
        ], size=10, color=MUTED)
        self._tb(slide, Inches(5.4), Inches(4.2), Inches(3.8), Inches(0.3),
                 "Uses SBB-XCUT-002 (Standalone)", size=10, bold=True, color=BLUE)

    # ------------------------------------------------------------------
    # SBB Swappability
    # ------------------------------------------------------------------

    def _make_swappability_slide(self, prs, data):
        slide = self._new_slide(prs, BG_CONTENT)
        self._tb(slide, Inches(0.7), Inches(0.35), Inches(8.0), Inches(0.6),
                 "SBB Swappability in Practice", size=24, bold=True)
        self._tb(slide, Inches(0.7), Inches(0.85), Inches(8.5), Inches(0.35),
                 "Multiple SBBs can satisfy the same ABB contract. When requirements change, swap the SBB — not the architecture.",
                 size=11, color=MUTED)

        # Build ABB → SBB mappings from data
        patterns = data["patterns"]
        abb_sbb_map = defaultdict(list)
        for p in patterns:
            if p.get("type") == "SBB":
                rels = p.get("relationships", [])
                for r in rels:
                    if r.get("type") == "IMPLEMENTS":
                        abb_sbb_map[r["target_id"]].append(p)

        # Show top ABBs with multiple SBBs
        entries = [(aid, sbbs) for aid, sbbs in abb_sbb_map.items() if len(sbbs) >= 2]
        entries.sort(key=lambda x: -len(x[1]))

        y_val = 1.5
        for abb_id, sbbs in entries[:3]:
            y = Inches(y_val)
            abb_name = self._name_map.get(abb_id, "")
            # ABB card background
            self._rect(slide, Inches(0.5), y, Inches(3.2), Inches(1.2), CARD_BG)
            self._rect(slide, Inches(0.5), y, Inches(0.06), Inches(1.2), BLUE)
            self._tb_multi(slide, Inches(0.8), y + Inches(0.2), Inches(2.7), Inches(0.8),
                           [abb_id, abb_name], size=12, color=WHITE, bold=True)
            # Arrow
            self._tb(slide, Inches(3.7), y, Inches(0.8), Inches(1.2),
                     "→", size=24, color=BLUE)
            # SBB rows
            sbb_y_val = y_val
            for sbb in sorted(sbbs, key=lambda s: s["id"])[:4]:
                # SBB row background
                self._rect(slide, Inches(4.8), Inches(sbb_y_val), Inches(4.7), Inches(0.35), ROW_BG)
                self._rect(slide, Inches(4.8), Inches(sbb_y_val), Inches(0.04), Inches(0.35), GREEN)
                self._tb(slide, Inches(5.05), Inches(sbb_y_val + 0.02), Inches(4.3), Inches(0.25),
                         f"{sbb['id']}  {sbb.get('name', '')}", size=10, color=MUTED)
                sbb_y_val += 0.4
            y_val += 1.7

    # ------------------------------------------------------------------
    # ABB Inventory
    # ------------------------------------------------------------------

    def _make_abb_inventory(self, prs, abbs, cat_map):
        slide = self._new_slide(prs, BG_CONTENT)
        self._tb(slide, Inches(0.7), Inches(0.35), Inches(8.0), Inches(0.5),
                 "Architecture Building Blocks", size=24, bold=True)
        self._tb(slide, Inches(0.7), Inches(0.8), Inches(8.0), Inches(0.3),
                 f"{len(abbs)} ABBs defining vendor-neutral capability contracts", size=12, color=MUTED)

        # Header bar
        self._rect(slide, Inches(0.5), Inches(1.25), Inches(9.0), Inches(0.35), BLUE)
        self._tb(slide, Inches(0.7), Inches(1.27), Inches(2.0), Inches(0.3),
                 "ID", size=10, bold=True)
        self._tb(slide, Inches(2.7), Inches(1.27), Inches(4.0), Inches(0.3),
                 "Name", size=10, bold=True)
        self._tb(slide, Inches(7.2), Inches(1.27), Inches(2.0), Inches(0.3),
                 "Category", size=10, bold=True)

        y_val = 1.65
        for i, abb in enumerate(sorted(abbs, key=lambda p: p["id"])):
            if y_val > 5.0:
                break
            y = Inches(y_val)
            # Alternating row backgrounds
            row_bg = CARD_BG if i % 2 == 0 else ROW_BG
            self._rect(slide, Inches(0.5), y, Inches(9.0), Inches(0.38), row_bg)
            cat_label = cat_map.get(abb.get("category", ""), abb.get("category", ""))
            self._tb(slide, Inches(0.7), y + Inches(0.04), Inches(2.0), Inches(0.25),
                     abb["id"], size=10, bold=True, color=BLUE)
            name_text = abb.get("name", "")
            tags = abb.get("tags", [])
            if tags:
                name_text += f"  [{', '.join(tags[:3])}]"
            self._tb(slide, Inches(2.7), y + Inches(0.04), Inches(4.0), Inches(0.25),
                     name_text, size=10)
            self._tb(slide, Inches(7.2), y + Inches(0.04), Inches(2.0), Inches(0.25),
                     cat_label, size=10, color=MUTED)
            y_val += 0.4

    # ------------------------------------------------------------------
    # SBB Inventory
    # ------------------------------------------------------------------

    def _make_sbb_inventory(self, prs, sbbs):
        slide = self._new_slide(prs, BG_CONTENT)
        self._tb(slide, Inches(0.7), Inches(0.35), Inches(8.0), Inches(0.5),
                 "Solution Building Blocks", size=24, bold=True)
        self._tb(slide, Inches(0.7), Inches(0.8), Inches(8.0), Inches(0.3),
                 f"{len(sbbs)} SBBs providing concrete, vendor-specific implementations", size=12, color=MUTED)

        sorted_sbbs = sorted(sbbs, key=lambda p: p["id"])
        half = (len(sorted_sbbs) + 1) // 2

        # Left column
        y_val = 1.2
        for i, sbb in enumerate(sorted_sbbs[:half]):
            if y_val > 5.0:
                break
            y = Inches(y_val)
            # Row background
            row_bg = CARD_BG if i % 2 == 0 else ROW_BG
            self._rect(slide, Inches(0.5), y, Inches(4.2), Inches(0.35), row_bg)
            self._tb(slide, Inches(0.7), y + Inches(0.03), Inches(1.5), Inches(0.25),
                     sbb["id"], size=8, bold=True, color=GREEN)
            self._tb(slide, Inches(2.1), y + Inches(0.03), Inches(2.5), Inches(0.25),
                     sbb.get("name", ""), size=9)
            y_val += 0.4

        # Right column
        y_val = 1.2
        for i, sbb in enumerate(sorted_sbbs[half:]):
            if y_val > 5.0:
                break
            y = Inches(y_val)
            # Row background
            row_bg = CARD_BG if i % 2 == 0 else ROW_BG
            self._rect(slide, Inches(5.3), y, Inches(4.2), Inches(0.35), row_bg)
            self._tb(slide, Inches(5.5), y + Inches(0.03), Inches(1.5), Inches(0.25),
                     sbb["id"], size=8, bold=True, color=GREEN)
            self._tb(slide, Inches(6.9), y + Inches(0.03), Inches(2.5), Inches(0.25),
                     sbb.get("name", ""), size=9)
            y_val += 0.4

    # ------------------------------------------------------------------
    # Closing Slide
    # ------------------------------------------------------------------

    def _make_closing_slide(self, prs, data):
        slide = self._new_slide(prs, BG_TITLE)

        # Top banner + accent line
        self._rect(slide, Inches(0), Inches(0), Inches(10), Inches(2.4), BG_DIVIDER)
        self._rect(slide, Inches(0), Inches(2.35), Inches(10), Inches(0.04), BLUE)

        self._tb(slide, Inches(0.7), Inches(0.5), Inches(8.0), Inches(0.6),
                 "Architecture Patterns", size=28, bold=True)
        # Export date with timezone
        from zoneinfo import ZoneInfo
        est_now = datetime.now(ZoneInfo("America/New_York"))
        export_date = est_now.strftime("%B %d, %Y  %I:%M %p EST")
        self._tb(slide, Inches(0.7), Inches(1.2), Inches(5.0), Inches(0.3),
                 export_date, size=11, color=MUTED)

        # Stats
        patterns = data["patterns"]
        ab_count = sum(1 for p in patterns if p.get("type") == "AB")
        abb_count = sum(1 for p in patterns if p.get("type") == "ABB")
        sbb_count = sum(1 for p in patterns if p.get("type") == "SBB")
        total = ab_count + abb_count + sbb_count

        stats = [
            (0.7, str(ab_count), "Architecture\nBlueprint"),
            (3.0, str(abb_count), "Architecture\nBuilding Blocks"),
            (5.3, str(sbb_count), "Solution\nBuilding Blocks"),
            (7.6, str(total), "Total\nPatterns"),
        ]
        for x, num, label in stats:
            self._tb(slide, Inches(x), Inches(2.8), Inches(2.0), Inches(0.7),
                     num, size=36, bold=True, color=BLUE)
            self._tb_multi(slide, Inches(x), Inches(3.5), Inches(2.0), Inches(0.5),
                           label.split("\n"), size=11, color=MUTED)

        # Next steps
        self._tb(slide, Inches(0.7), Inches(4.2), Inches(2.0), Inches(0.3),
                 "Next Steps", size=13, bold=True)
        self._tb(slide, Inches(0.7), Inches(4.5), Inches(8.5), Inches(0.6),
                 "Expand Level 4 patterns  ·  Add cost/performance benchmarks  ·  Map to enterprise use cases  ·  Governance integration",
                 size=10, color=MUTED)
